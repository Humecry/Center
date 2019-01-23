#! /usr/bin/env python3
# -*- coding:utf-8 -*-

# @Date    : 2018-08-03 16:45:33
# @Author  : Hume (102734075@qq.com)
# @Link    : https://humecry.wordpress.com/
# @Version : 1.0
# @Description: 公共函数

import sys
import time
import datetime
from datetime import date
from dateutil.relativedelta import relativedelta
import pyodbc
import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.shapes.table import CT_Table
from functools import partial
# 引入配置文件
from conf import *

# 统计时间
today = date.today()
# today = datetime.datetime.strptime("2018-09-04", "%Y-%m-%d")

'''
months int 值为-1则为上月日期，为-2则为上上月日期
type str 值为first则返回第一天，值为end则返回最后一天
'''
def lastMonthDate(months, type='first'):
	if type == 'first':
		lastMonth = today - relativedelta(months=-months)
		# 第一天
		lastMonthFirstDate = date(lastMonth.year, lastMonth.month, 1)
		lastMonthFirstDateStr = lastMonthFirstDate.strftime("%Y%m%d")
		return lastMonthFirstDateStr
	elif type == 'end':
		lastMonthAfter = today - relativedelta(months=-months-1)
		# 最后一天
		lastMonthEndDate = date(lastMonthAfter.year, lastMonthAfter.month, 1) - relativedelta(days=1)
		lastMonthEndDateStr = lastMonthEndDate.strftime("%Y%m%d")
		return lastMonthEndDateStr
	else:
		return None

# 设置PPT关于文本的格式
def ppt_layout(cell, value=None, align=PP_ALIGN.CENTER, align_v=MSO_ANCHOR.MIDDLE, fontName='微软雅黑', bold = False, color=(0, 0, 0)):
	if value:
		cell.text = value
	cell.text_frame.paragraphs[0].alignment = align
	cell.vertical_anchor = align_v
	cell.text_frame.paragraphs[0].font.name = fontName
	cell.text_frame.paragraphs[0].font.bold = bold
	cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(color[0], color[1], color[2])

# PPT表格上下合并单元格
def mergeCellsVertically(table, start_row_idx, end_row_idx, col_idx):
	row_count = end_row_idx - start_row_idx + 1
	column_cells = [r.cells[col_idx] for r in table.rows][start_row_idx:]
	column_cells[0]._tc.set('rowSpan', str(row_count))

# PPT表格左右合并单元格
def mergeCellsHorizontally(table, start_col_idx, end_col_idx, row_idx):
	col_count = end_col_idx - start_col_idx + 1
	row_cells = [c for c in table.rows[row_idx].cells][start_col_idx:end_col_idx]
	row_cells[0]._tc.set('gridSpan', str(col_count))

# PPT表格大范围合并单元格
def mergeCells(table, start_row_idx, end_row_idx, start_col_idx, end_col_idx):
	for col_idx in range(start_col_idx,end_col_idx+1):
		mergeCellsVertically(table, start_row_idx, end_row_idx, col_idx)
	for row_idx in range(start_row_idx, end_row_idx+1):
		mergeCellsHorizontally(table, start_col_idx, end_col_idx, row_idx)

"""
显示处理进度的类
调用该类相关函数即可实现处理进度的显示
"""
class ShowProcess():
    i = 0 # 当前的处理进度
    max_steps = 0 # 总共需要处理的次数
    max_arrow = 50 #进度条的长度
    infoDone = 'done'

    # 初始化函数，需要知道总共的处理次数
    def __init__(self, max_steps, infoDone = 'Done'):
        self.max_steps = max_steps
        self.i = 0
        self.infoDone = infoDone

    # 显示函数，根据当前的处理进度i显示进度
    # 效果为[>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>]100.00%
    def show_process(self, i=None):
        if i is not None:
            self.i = i
        else:
            self.i += 1
        num_arrow = int(self.i * self.max_arrow / self.max_steps) #计算显示多少个'>'
        num_line = self.max_arrow - num_arrow #计算显示多少个'-'
        percent = self.i * 100.0 / self.max_steps #计算完成进度，格式为xx.xx%
        process_bar = '[' + '>' * num_arrow + '-' * num_line + ']'\
                      + '%.2f' % percent + '%' + '\r' #带输出的字符串，'\r'表示不换行回到最左边
        sys.stdout.write(process_bar) #这两句打印字符到终端
        sys.stdout.flush()
        if self.i >= self.max_steps:
            self.close()

    def close(self):
        print('')
        print(self.infoDone)
        self.i = 0